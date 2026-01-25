#!/usr/bin/env python3
"""
MongoDB CSFLE & Queryable Encryption - Ultra-Detailed SA Enablement Presentation
Generates a 25-slide high-fidelity PPTX with deep technical content and speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# MongoDB brand colors
MONGODB_GREEN = RGBColor(0, 104, 74)   # #00684A
MONGODB_DARK = RGBColor(33, 49, 60)    # #21313C
MONGODB_LIGHT = RGBColor(249, 251, 250) # #F9FBFA

def add_title_slide(prs, title, subtitle="", footer=""):
    """Add a professional title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background accent (optional, keeping it clean for now)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = MONGODB_DARK
        p.alignment = PP_ALIGN.CENTER
        
    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        p = footer_box.text_frame.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = MONGODB_DARK
        p.alignment = PP_ALIGN.CENTER
        
    return slide

def add_content_slide(prs, title, bullets, speaker_notes=""):
    """Add a content slide with bullet points and MongoDB branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = MONGODB_GREEN
    line.line.fill.background()
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(bullet)
        p.font.size = Pt(20)
        p.font.color.rgb = MONGODB_DARK
        p.space_after = Pt(10)
        
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes
        
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, speaker_notes=""):
    """Add a two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    # Left Column
    lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.5))
    p = lt_box.text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    l_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(4.5))
    tf_l = l_box.text_frame
    tf_l.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = str(b)
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        
    # Right Column
    rt_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    p = rt_box.text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    r_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(4.5))
    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.text = str(b)
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes
        
    return slide

def add_code_slide(prs, title, code, speaker_notes=""):
    """Add a dark code slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(30,30,30)
    
    c_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(4.6))
    tf = c_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(230,230,230)
    
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes
    return slide

def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    add_title_slide(prs, "MongoDB Client-Side Field Level Encryption & Queryable Encryption", 
                    "SA Technical Enablement Deep-Dive", 
                    "45-minute presentation + 3 hands-on labs")
    prs.slides[0].notes_slide.notes_text_frame.text = "Welcome to the CSFLE & Queryable Encryption deep-dive. This session is designed for senior SAs who need to handle complex 'Day 2' architectural discussions with customers. By the end, you'll explain HOW QE works, design multi-cloud architectures, and handle GDPR/HIPAA with confidence."

    # 2. Agenda
    add_content_slide(prs, "Agenda", 
                      ["0-5 min: The 'Why' & Compliance Hook", 
                       "5-15 min: Cryptographic Fundamentals & Internals",
                       "15-25 min: GDPR & Multi-Cloud Patterns",
                       "25-35 min: CSFLE vs QE Architecture Differences",
                       "35-45 min: Competitive 'Kill' Tracks"],
                      "We start with the business case, then go deep into cryptographic internals—the 'how it works' that sets SAs apart. GDPR patterns are critical for Europe. Comparison is where customers get confused. Finally, competitive positioning.")

    # 3. Cost of Getting It Wrong
    add_content_slide(prs, "The Cost of Getting It Wrong",
                      ["€4.4B+ in GDPR fines issued as of 2024",
                       "$1.3M average HIPAA penalty per violation",
                       "GDPR Art 32: 'Encryption of personal data' required",
                       "GDPR Art 17: 'Right to be Forgotten' requires erasure",
                       "PCI-DSS Req 3: Encrypt PAN, manage crypto keys",
                       "SOX: Audit trail + access controls required"],
                      "These aren't theoretical risk—€4.4B in fines... Art 32 requires encryption... Art 17 requires full deletion—we'll show 'crypto-shredding' for this. HIPAA requires ePHI protection at rest/transit. PCI-DSS needs encryption + search.")

    # 4. Use Cases by Industry
    add_content_slide(prs, "Industry Use Cases",
                      ["Healthcare: ePHI, Insurance IDs, Prescriptions (HIPAA)",
                       "Financial: Account numbers, SSN, Tax IDs (PCI-DSS, SOX)",
                       "Gaming/Social: Payments, Chat logs, Identity (COPPA, GDPR)",
                       "Differentiator: Selective encryption + functional queryability"],
                      "Pattern: Selective encryption while maintaining app function. Health: encrypt SSN but leave department searchable. Finance: Range queries on encrypted transaction amounts for fraud detection.")

    # 5. CSFLE vs QE Overview
    add_two_column_slide(prs, "CSFLE vs Queryable Encryption",
                         "CSFLE (4.2+)",
                         ["Deterministic/Random modes", "Equality queries only", "Mature/Battle-tested", "Lower overhead", "Share 1 DEK across fields"],
                         "QE (7.0+)",
                         ["Randomized encryption only", "Equality AND Range queries", "EMM-based crypto breakthrough", "2-3x storage overhead", "REQUIRES separate DEK per field"],
                         "Fundamental choice: CSFLE is mature, simpler. QE enables Range but requires a SEPARATE DEK for EACH encrypted field because of metadata binding. Crucial distinction.")

    # 6. Envelope Encryption Architecture
    add_content_slide(prs, "Envelope Encryption Architecture",
                      ["KMS Provider (AWS/Azure/GCP/KMIP) - Holds CMK",
                       "MongoDB Key Vault - Holds Encrypted DEKs",
                       "Application Data - Encrypted by DEK (BSON Subtype 6)",
                       "Flow: Client retrieves DEK -> Decrypt via KMS -> Encrypt Data",
                       "Benefit: CMK rotation is cheap (only re-encrypt DEKs)"],
                      "Three layers: CMK in KMS, DEK in Key Vault, Data in BSON. Client retrives encrypted DEK, KMS decrypts using CMK, client encrypts data. Defense in depth—CMK never leaves KMS.")

    # 7. Structured Encryption & EMMs
    add_content_slide(prs, "QE: Structured Encryption & EMMs",
                      ["Encrypted Multi-Maps (EMMs) enable server-side logic",
                       "Client generates 'tokens' encoding order relations",
                       "Server tests tokens ($gt, $lt, $eq) without plaintext",
                       "Private Querying: Access patterns redacted from logs",
                       "Security: Randomized encryption still searchable"],
                      "How can server do $gt on ciphertext? EMMs. Client generates 'range tokens' encoding relations. Server tests tokens without knowing values. Logs don't even show WHICH field is queried.")

    # 8. .esc and .ecoc Collections
    add_content_slide(prs, "QE Metadata Collections",
                      [".esc (System Catalog): Metadata mapping fields to DEKs",
                       ".ecoc (Context Cache): Stores query tokens",
                       "Impact: Sizing must account for these collections",
                       "Operations: Monthly compaction recommended",
                       "Backups: Must include enxcol_ collections"],
                      "QE creates enxcol_ collections. .esc maps fields to DEKs. .ecoc stores search tokens—this grows and needs monthly compaction (online operation). Factor this into storage planning.")

    # 9. Storage Overhead Challenge
    add_content_slide(prs, "Challenge: Storage Overhead",
                      ["Q: What is the overhead for a QE Range field?",
                       "A) 1.2x", "B) 1.5x", "C) 2-3x", "D) 5x+",
                       "Answer: C (2-3x overhead)",
                       "Why? Multiple tokens for range capability per field",
                       "Tuning: sparsity and min/max settings impact size"],
                      "Range queries require multiple tokens per value for math relations. Result is 2-3x overhead. Tuning min/max precisely saves space. Sparsity (1-4) density impacts query precision vs size.")

    # 10. Key Hierarchy Visualization
    add_content_slide(prs, "The Key Protection Stack",
                      ["☁️ KMS Provider: Enterprise Management",
                       "🔐 CMK (Master): Protects the DEKs",
                       "🔑 DEKs (Data): Stored encrypted, 1 per field (QE)",
                       "📄 Field Data: BSON Subtype 6",
                       "Zero-Downtime Rotation: rewrapManyDataKey()"],
                      "KMS -> CMK -> DEK -> Data. Rotation is fast—only re-encrypt small DEK documents. Actual data remains unchanged. compliance mandatory but operationally low impact.")

    # 11. Right to Erasure / GDPR
    add_content_slide(prs, "GDPR: Crypto-Shredding at Scale",
                      ["Problem: backups/logs contain user data after deletion",
                       "Pattern: 'One DEK per User'",
                       "Action: Delete user's DEK from Key Vault",
                       "Result: All user data (including backups) is unrecoverable",
                       "Compliance: Statisfies Art 17 'Right to be Forgotten'"],
                      "Backups make deletion hard. crypto-shredding is the answer. 1 user = 1 DEK. Delete DEK = data is garbage. Lab 3 will implement this. Note: cross-user scaling needs special design.")

    # 12. Multi-Cloud KMS Architecture
    add_two_column_slide(prs, "Multi-Cloud KMS Architecture",
                         "BYOK (Bring Your Own Key)",
                         ["Import key material to Cloud KMS", "Full lifecycle control", "Regulation custody compliant", "Higher ops complexity"],
                         "Managed Identity",
                         ["Cloud provider generates/rotates", "Low overhead / Native integration", "Simpler ops", "Sufficient for most audits"],
                         "AWS KMS (IAM), Azure Key Vault (Managed Identity), GCP Cloud KMS, KMIP. Tip: use different providers for Prod vs DR for true resilience.")

    # 13. Automatic vs Explicit Encryption
    add_two_column_slide(prs, "Integration Strategies",
                         "🤖 Automatic", ["Define encryptedFields map", "Driver intercepts & encrypts", "No app code changes", "Required for QE Range"],
                         "👩‍💻 Explicit", ["Manual encrypt/decrypt calls", "Full control over logic", "Retrofitting legacy apps", "Conditional logic support"],
                         "Automatic is best for new apps and QE. Explicit offers fine-grained control when automatic rules are too rigid for business logic.")

    # 14. QE: The encryptedFields Map
    add_code_slide(prs, "QE Configuration Map",
                   """const encryptedFields = {
  fields: [
    { path: "ssn", bsonType: "string", queries: { queryType: "equality" } },
    { path: "salary", bsonType: "int",
      queries: { queryType: "range", min: 0, max: 1000000,
                 sparsity: 2, contention: 4 } }
  ]
};""",
                   "min/max are REQUIRED for range—tighter bounds = fewer tokens. sparsity (density) and contention (write concurrency) are tunable knobs for performance vs security.")

    # 15. Query Support Matrix
    add_content_slide(prs, "Query Support Matrix",
                      ["✅ $eq, $ne, $in - All encrypted modes",
                       "✅ $gt, $lt, $gte, $lte - QE Range only",
                       "❌ $regex, $text search - Unsupported",
                       "❌ Sorting - Unsupported (leaks order)",
                       "❌ Aggregations ($group, $sum) - Unsupported",
                       "Contention Factor: balances throughput vs frequency analysts"],
                      "Search limits: no regex or text search. No server-side sorting. No aggregations—do these on decrypted data in the app or use separate anonymized stores.")

    # 16. Honest Limitations & Workarounds
    add_two_column_slide(prs, "Honest Limitations",
                         "❌ No-go Zones", ["Sorting ciphertext", "Regex/Partial matches", "Server-side aggregations", "Atlas Search on encrypted fields"],
                         "✅ SA Workarounds", ["Client-side sort after decrypt", "Tokenized search fields", "Anonymized aggregation stores", "Separate analytics read models"],
                         "Better to be honest upfront. If they need heavy sorting or Atlas Search on sensitive data, encryption must be designed carefully around those needs.")

    # 17. Performance & Operations
    add_content_slide(prs, "Ops & Performance Impact",
                      ["Storage: 2-3x (Range), 1.5x (Equality)",
                       "Write Latency: ~10% insert overhead",
                       "Writes: 1 insert -> multiple internal writes",
                       "Compaction: Mandatory monthly online cleanup",
                       "Maintenance: Include enxcol_ collections in backups"],
                      "Storage is the main cost. Range = 2-3x factor. Latency from transactional internal writes to .esc/.ecoc. Compaction is online, no downtime.")

    # 18. KMS & Key Rotation
    add_code_slide(prs, "Rotating Keys (Zero-Downtime)",
                   """await clientEncryption.rewrapManyDataKey({}, {
  provider: "aws",
  masterKey: { key: "arn:aws:kms...NEW-CMK" }
});""",
                   "rewrapManyDataKey re-encrypts DEKs with the new Master Master Key. Actual documents are never touched. Complies with annual rotation mandates easily.")

    # 19. Regulatory Alignment
    add_content_slide(prs, "Regulatory Compliance Mapping",
                      ["GDPR Art 32: Field-level encryption",
                       "GDPR Art 17: 1 DEK per user (Erasure)",
                       "HIPAA: ePHI as BSON Subtype 6",
                       "PCI-DSS Req 3: Searchable encrypted PAN",
                       "Audit: KMS access logs for auditors"],
                      "auditor artifacts: BSON Subtype 6 docs, KMS CloudTrail logs, rewrapManyDataKey rotation schedules. Proven data protection in use.")

    # 20. Competitive Kill Tracks
    add_content_slide(prs, "Competitive Differentiation",
                      ["vs Oracle TDE: TDE decrypts in memory (DBAs see data)",
                       "vs Cosmos DB: Cosmos is deterministic-only (leaks patterns)",
                       "vs PostgreSQL: Server-side, no KMS integration",
                       "MongoDB: True zero-trust (server never sees plaintext)",
                       "Differentiator: Random encryption + Range queries"],
                      "TDE is just disk-at-rest. QE is encryption-in-use. Cosmos lacks range queries on encrypted data. MongoDB is the leader in usable zero-trust.")

    # 21. Customer Discovery Questions
    add_content_slide(prs, "SA Discovery Questions",
                      ["Who should NOT see this data? (The killer question)",
                       "Is it acceptable for DBAs to see customer SSNs?",
                       "Do you have 'Right to be Forgotten' requirements?",
                       "What is your key management security model?",
                       "Do you need to search encrypted data?"],
                      "Identifying the 'DBA Admin' threat usually unlocks the budget for client-side encryption. Ask specifically if admins should see sensitive card IDs.")

    # 22. Objection Handling
    add_content_slide(prs, "Objection Handling",
                      ["'We have TDE' -> TDE leaks to DBAs; we don't.",
                       "'It's too slow' -> security/performance trade-off; encrypt less.",
                       "'Can't search' -> QE enables Equality and Range.",
                       "'Too complex' -> Automatic mode needs no app changes.",
                       "'Compliance loves TDE' -> mention 'unauthorized access' clauses."],
                      "Frame it as Risk vs Performance. 10% latency vs $4.5M average breach cost. Use automatic encryption to rebut 'Complexity' objections.")

    # 23. Anti-Patterns
    add_content_slide(prs, "Encryption Anti-Patterns",
                      ["Don't encrypt fields for Sorting",
                       "Don't encrypt fields for Heavy Aggregations",
                       "Don't encrypt EVERYTHING (security overkill/perf hit)",
                       "Don't use Local Key in production",
                       "Don't skip compaction for active users"],
                      "Qualify early. If they MUST group by the encrypted field on the server, QE is the wrong fit. Design the schema to encrypt only the PII data.")

    # 24. Hands-on Labs Overview
    add_content_slide(prs, "Hands-on Training (3 x 34 min)",
                      ["Lab 1: CSFLE Fundamentals & Troubleshooting",
                       "Lab 2: QE Internals, .esc/.ecoc, and Sizing",
                       "Lab 3: GDPR Compliance & Crypto-Shredding",
                       "Scenario: AWS KMS integration (Enterprise focus)"],
                      "Lab 1 fixes common local vs atlas pathing errors. Lab 2 goes deep into metadata storage. Lab 3 builds the Erasure workflow.")

    # 25. Key Takeaways
    add_content_slide(prs, "Key Takeaways",
                      ["QE = Randomized + Range Queries (the big one)",
                       "2-3x storage for Range - Capacity Plan now",
                       "Mandatory DEK per field (QE metadata binding)",
                       "Online Compaction for ongoing performance",
                       "Crypto-shredding satisfies GDPR Art 17",
                       "docs.mongodb.com/qe"],
                      "Final checklist: 1. QE Range factor 2. Compaction schedule 3. One DEK per field. 4. Crypto-shredding strategy. Move to Lab 1.")

    prs.save("MongoDB_Advanced_Security_SA_DeepDive.pptx")
    print("✅ PowerPoint generated successfully.")

if __name__ == "__main__":
    generate_presentation()
