#!/usr/bin/env python3
"""
MongoDB CSFLE & Queryable Encryption SA Enablement Presentation Generator
Generates a 25-slide PowerPoint presentation with speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# MongoDB brand colors
MONGODB_GREEN = RGBColor(0, 104, 74)  # #00684A
MONGODB_DARK = RGBColor(33, 49, 60)   # #21313C
MONGODB_LIGHT = RGBColor(249, 251, 250)  # #F9FBFA

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.75))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = MONGODB_DARK
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, speaker_notes=""):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    # Green underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = MONGODB_GREEN
    line.line.fill.background()
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = MONGODB_DARK
        p.space_after = Pt(12)
    
    # Speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, speaker_notes=""):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MONGODB_DARK
    
    # Left column bullets
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = MONGODB_DARK
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MONGODB_DARK
    
    # Right column bullets
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = MONGODB_DARK
        p.space_after = Pt(8)
    
    # Speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes
    
    return slide

def add_code_slide(prs, title, code, speaker_notes=""):
    """Add a slide with code block."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MONGODB_GREEN
    
    # Code background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
    
    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(4.6))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(12)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(220, 220, 220)
    
    # Speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes
    
    return slide

def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # =====================================================================
    # SLIDE 1: Title Slide
    # =====================================================================
    add_title_slide(prs, 
        "MongoDB CSFLE & Queryable Encryption",
        "SA Technical Enablement Deep-Dive"
    )
    
    # =====================================================================
    # SLIDE 2: Agenda
    # =====================================================================
    add_content_slide(prs, "Today's Journey", [
        "0-5 min: The \"Why\" & Compliance Hook",
        "5-15 min: Cryptographic Fundamentals & Internals",
        "15-25 min: GDPR & Multi-Cloud Patterns",
        "25-35 min: CSFLE vs QE Architecture Differences",
        "35-45 min: Competitive \"Kill\" Tracks"
    ], """This is a technical deep-dive for senior SAs. By the end:
- You'll explain HOW QE works under the hood
- Design multi-cloud security architectures
- Handle GDPR/HIPAA compliance questions with confidence
- Differentiate from Oracle TDE and Cosmos DB""")
    
    # =====================================================================
    # SLIDE 3: The Compliance Landscape
    # =====================================================================
    add_content_slide(prs, "The Cost of Getting It Wrong", [
        "€4.4B+ in GDPR fines to date",
        "Average HIPAA penalty: $1.3M",
        "GDPR Art. 32: \"Encryption of personal data\"",
        "HIPAA: ePHI must be encrypted at rest AND in transit",
        "PCI-DSS: Encrypt cardholder data, manage crypto keys",
        "100% of major regulations now require encryption"
    ], """GDPR Article 32 specifically calls out encryption.
Article 17 - the \"Right to be Forgotten\" - requires you to delete all user data.
CSFLE enables \"crypto-shredding\" to address this.
HIPAA requires encryption of electronic Protected Health Information (ePHI).
PCI-DSS is interesting because you need to encrypt cardholder data BUT also need to search it.""")
    
    # =====================================================================
    # SLIDE 4: Use Cases by Industry
    # =====================================================================
    add_content_slide(prs, "Where CSFLE & QE Shine", [
        "Healthcare: Patient records (ePHI), Insurance IDs, Prescriptions → HIPAA",
        "Financial Services: Account numbers, SSN/Tax IDs, Transaction amounts → PCI-DSS, SOX",
        "Gaming & Social: Payment details, Chat logs (minors), Location → COPPA, GDPR",
        "Government: Classified data, PII, Biometrics → FedRAMP, ITAR"
    ], """Pattern: Selective encryption of sensitive fields while maintaining application functionality.
Healthcare: Encrypt patient_name, SSN, diagnosis codes - leave timestamps queryable.
Financial: QE enables range queries on encrypted transaction amounts for fraud detection.""")
    
    # =====================================================================
    # SLIDE 5: CSFLE vs Queryable Encryption Overview
    # =====================================================================
    add_two_column_slide(prs, "Two Solutions, Different Trade-offs",
        "CSFLE (MongoDB 4.2+)",
        [
            "Deterministic OR Random encryption",
            "Equality queries on deterministic fields",
            "Mature, battle-tested",
            "Lower overhead",
            "Can share DEK across multiple fields"
        ],
        "Queryable Encryption (7.0+)",
        [
            "Always random encryption",
            "Equality AND Range queries",
            "Latest innovation",
            "2-3x storage overhead",
            "Requires separate DEK per field"
        ],
        """CRITICAL: QE requires a SEPARATE DEK for EACH encrypted field due to metadata binding.
In CSFLE, you might use one DEK for all sensitive fields.
In QE, each field needs its own DEK because of how the metadata collections work.""")
    
    # =====================================================================
    # SLIDE 6: Envelope Encryption Architecture
    # =====================================================================
    add_content_slide(prs, "The Envelope Encryption Model", [
        "CMK (Customer Master Key) → Lives in KMS, NEVER leaves",
        "DEK (Data Encryption Key) → Stored encrypted in Key Vault",
        "Data → Encrypted with DEK, stored as BSON Subtype 6",
        "Flow: Client → KMS (decrypt DEK) → Use DEK → Encrypt data",
        "CMK rotation only requires re-encrypting DEKs, not data"
    ], """Three layers of protection:
1. CMK in KMS (AWS, Azure, GCP, KMIP) - your master key
2. DEKs in Key Vault collection - encrypted by CMK
3. Data as BSON Subtype 6 in your documents

Why envelope encryption?
- CMK rotation is cheap (only re-encrypt DEKs)
- CMK never touches your infrastructure
- Defense in depth""")
    
    # =====================================================================
    # SLIDE 7: Structured Encryption & EMMs
    # =====================================================================
    add_content_slide(prs, "How QE Enables Range Queries on Ciphertext", [
        "Encrypted Multi-Maps (EMMs) enable server-side computation",
        "Client generates encrypted tokens at insert time",
        "Tokens encode order relationships mathematically",
        "Server can test $gt, $lt, $eq without seeing plaintext",
        "Private Querying: Query patterns redacted from logs"
    ], """This is the deep technical content that sets you apart.

Conceptually:
1. Client generates special \"range tokens\" at encryption time
2. These tokens encode order relationships mathematically
3. Server can test token relationships without knowing values
4. Result: Server finds \"salary > 50000\" without ever seeing the values

Private Querying is a bonus - a DBA can't even tell WHICH field you're querying.""")
    
    # =====================================================================
    # SLIDE 8: .esc and .ecoc Collections
    # =====================================================================
    add_content_slide(prs, "QE Internal Collections", [
        ".esc (Encrypted State Collection / System Catalog):",
        "  - Named: enxcol_.<collection>.esc",
        "  - Stores metadata about encrypted fields",
        "  - Maps each field to its DEK",
        ".ecoc (Encrypted Compaction Collection / Context Cache):",
        "  - Named: enxcol_.<collection>.ecoc", 
        "  - Stores query tokens generated during inserts",
        "  - Grows with each insert - needs periodic compaction"
    ], """Why this matters for customers:
1. Storage planning - these collections add overhead
2. Operational procedures - compaction should be scheduled monthly
3. Backup considerations - these collections must be backed up too

In Lab 2, you'll explore these collections in Compass.""")
    
    # =====================================================================
    # SLIDE 9: Storage Factor Challenge
    # =====================================================================
    add_content_slide(prs, "Quick Challenge: Storage Overhead", [
        "Question: What's the storage overhead for a Range-indexed encrypted field?",
        "",
        "A) 1.2x",
        "B) 1.5x", 
        "C) 2-3x",
        "D) 5x+",
        "",
        "Answer: C) 2-3x for Range-indexed fields"
    ], """Why 2-3x?
- Each value generates multiple tokens for range capability
- Tokens stored in .esc and .ecoc collections
- More granular ranges = more tokens

Factors affecting overhead:
- Number of encrypted fields
- Range index sparsity setting (min/max)
- Cardinality of values
- Query precision requirements""")
    
    # =====================================================================
    # SLIDE 10: Key Hierarchy Visualization
    # =====================================================================
    add_content_slide(prs, "The Key Protection Stack", [
        "KMS Provider (AWS/Azure/GCP/KMIP)",
        "    ↓ protects",
        "CMK (Customer Master Key)",
        "    ↓ encrypts",
        "DEK(s) (Data Encryption Keys)",
        "    ↓ encrypts",
        "Field Data (BSON Subtype 6)"
    ], """Benefits of this architecture:
- CMK rotation is cheap (only re-encrypt DEKs, not all data)
- Separation of concerns (KMS team manages CMK, app team uses DEKs)
- Audit trail at every layer
- Cloud-agnostic (can switch KMS providers)

Key rotation workflow using rewrapManyDataKey():
1. Create new CMK in KMS
2. Call rewrapManyDataKey() to re-encrypt all DEKs
3. Retire old CMK - no data re-encryption needed!""")
    
    # =====================================================================
    # SLIDE 11: Right to Erasure Pattern
    # =====================================================================
    add_content_slide(prs, "GDPR Article 17: Crypto-Shredding", [
        "The Problem: User requests deletion, but data is in backups, logs, everywhere",
        "The Solution: One DEK Per User pattern",
        "",
        "1. User signs up → Generate unique DEK for them",
        "2. All their sensitive data encrypted with THEIR DEK",
        "3. User requests erasure → Delete their DEK from Key Vault",
        "4. Result: All their data becomes cryptographically indecipherable",
        "   (Including backups - backups have encrypted data)"
    ], """This is Lab 3 - implementing this pattern hands-on.

Limitation to discuss: Cross-user queries become complex.
If you need to query across users (analytics), you may need a separate aggregation approach.

This satisfies GDPR Art. 17 completely - the data is mathematically unrecoverable.""")
    
    # =====================================================================
    # SLIDE 12: Multi-Cloud KMS Architecture
    # =====================================================================
    add_two_column_slide(prs, "Managing CMKs Across Providers",
        "BYOK (Bring Your Own Key)",
        [
            "Customer owns and manages CMK",
            "Import existing keys to cloud KMS",
            "Full control over key lifecycle",
            "Required for some regulations"
        ],
        "Managed Identity",
        [
            "Cloud provider generates CMK",
            "Automatic rotation available",
            "Simpler operational model",
            "Sufficient for most use cases"
        ],
        """Provider options:
- AWS KMS: Most common, IAM-based auth
- Azure Key Vault: Managed Identity is excellent
- GCP Cloud KMS: Service account auth
- KMIP: For on-prem HSMs (Thales, Gemalto)

Architecture tip: You can use DIFFERENT KMS providers for different DEKs.
- Production data with AWS KMS
- DR data with Azure Key Vault
- Provides true multi-cloud resilience""")
    
    # =====================================================================
    # SLIDE 13: Automatic vs Explicit Encryption
    # =====================================================================
    add_two_column_slide(prs, "Choose Your Integration Path",
        "🤖 Automatic Encryption",
        [
            "Define encryptedFields on collection",
            "Driver intercepts all operations",
            "App code doesn't change",
            "Best for new applications",
            "Uses schemaMap or encryptedFields"
        ],
        "👩‍💻 Explicit Encryption",
        [
            "Call encrypt/decrypt methods directly",
            "Full control over when/what",
            "Can encrypt conditionally",
            "Best for retrofitting legacy apps",
            "Uses clientEncryption.encrypt()"
        ],
        """In practice, many customers use both:
- Automatic for the main path
- Explicit for edge cases or conditional encryption

Both use the same underlying crypto - it's about developer experience.""")
    
    # =====================================================================
    # SLIDE 14: Queryable Encryption Configuration
    # =====================================================================
    code = """// Modern QE Configuration (MongoDB 7.0+)
const encryptedFields = {
  fields: [
    {
      path: "ssn",
      bsonType: "string",
      queries: { queryType: "equality" }
    },
    {
      path: "salary",
      bsonType: "int",
      queries: { 
        queryType: "range",
        min: 0, 
        max: 1000000,
        sparsity: 2,       // Higher = fewer tokens
        contention: 4      // Balances throughput vs security
      }
    }
  ]
};

await db.createCollection("employees", { encryptedFields });"""
    add_code_slide(prs, "QE Configuration: encryptedFields", code,
        """Key parameters for Range queries:
- min/max: Bounds for token generation (required for Range)
- sparsity: Higher values = fewer tokens, lower precision
- contention: 0-16, balances write throughput vs frequency analysis protection

Important: Each field gets its own DEK automatically.""")
    
    # =====================================================================
    # SLIDE 15: Query Operator Support Matrix
    # =====================================================================
    add_content_slide(prs, "Query Operator Support Matrix", [
        "✅ $eq - All modes",
        "✅ $ne - All modes", 
        "✅ $gt, $gte, $lt, $lte - QE Range only",
        "✅ $in - All modes",
        "❌ $regex - Not supported",
        "❌ $text search - Not supported",
        "❌ Sorting - Would leak order information",
        "❌ $group/$sum - Aggregations not supported"
    ], """What you CANNOT do on encrypted fields:
- Text search ($text, $regex)
- Sorting (would leak order information)
- Aggregations ($group, $sum, $avg)

Schema design implication:
- Keep sortable/aggregatable fields unencrypted
- Or aggregate on the client side after decryption
- Consider separate analytics store for reporting""")
    
    # =====================================================================
    # SLIDE 16: Honest Limitations
    # =====================================================================
    add_two_column_slide(prs, "What You Can't Do (Be Honest)",
        "❌ Not Supported",
        [
            "Sorting on encrypted fields",
            "Full-text search ($text)",
            "Regex matching ($regex)",
            "Aggregation operators",
            "Array operations on encrypted arrays",
            "Computed fields using encrypted values"
        ],
        "✅ Workarounds",
        [
            "Client-side sorting after decryption",
            "Tokenized search on separate field",
            "Pre-computed aggregates (anonymized)",
            "Separate analytics collection",
            "Encrypt array as whole document"
        ],
        """Be honest with customers about limitations.
Better they hear it from you than discover it in production.

Best Practice: Design schema so sorting, searching, and aggregation
target non-sensitive fields. Encrypt only what needs protection.""")
    
    # =====================================================================
    # SLIDE 17: Performance Considerations
    # =====================================================================
    add_content_slide(prs, "Performance & Operational Impact", [
        "Storage: 2-3x overhead for Range fields, ~1.5x for Equality",
        "Write Latency: ~10% increase for encrypted inserts",
        "Each insert = writes to main doc + .esc + .ecoc + range tokens",
        "Compaction: .ecoc grows with inserts, compact monthly",
        "Compaction is online - no downtime required"
    ], """Performance facts:
- Storage: Plan for 2.5x factor for Range fields
- Writes: The overhead is from multiple internal writes (atomic)
- Compaction: Schedule monthly for active collections

Sizing guidance:
- Start with 2.5x storage factor for Range fields
- Monitor and adjust based on actual usage
- Consider separate storage tier for encrypted collections""")
    
    # =====================================================================
    # SLIDE 18: KMS Providers & Key Rotation
    # =====================================================================
    add_content_slide(prs, "Key Rotation with rewrapManyDataKey", [
        "Supported KMS Providers:",
        "  - AWS KMS (most common)",
        "  - Azure Key Vault",
        "  - GCP Cloud KMS",
        "  - KMIP (Thales, etc.)",
        "  - Local Key (DEV ONLY)",
        "",
        "Key Rotation is FAST:",
        "  - Only re-encrypts DEKs (small documents)",
        "  - Actual data remains unchanged",
        "  - 1M documents? Still just re-encrypting a few DEKs"
    ], """Key rotation with rewrapManyDataKey():
1. Create new CMK in your KMS
2. Call rewrapManyDataKey() - re-encrypts all DEKs
3. Retire old CMK after verification

Best practices:
- Rotate CMK annually (or per compliance requirements)
- Test rotation in staging first
- Keep old CMK around briefly for rollback
- Automate rotation in CI/CD pipeline""")
    
    # =====================================================================
    # SLIDE 19: Regulatory Alignment
    # =====================================================================
    add_content_slide(prs, "Compliance Coverage Matrix", [
        "GDPR Art. 32 → CSFLE/QE field-level encryption",
        "GDPR Art. 17 → 1 DEK per user = crypto-shredding",
        "HIPAA → CSFLE on PHI fields (BSON Subtype 6)",
        "PCI-DSS Req 3 → QE for searchable encrypted PAN",
        "SOX → KMS audit logs (CloudTrail, Azure Monitor)",
        "All → RBAC + encryption = defense in depth"
    ], """Auditor conversations:
- \"Show me the encryption\" → BSON Subtype 6 in documents
- \"Show me key management\" → KMS audit logs
- \"Show me access controls\" → RBAC + field-level separation

Key message: Encryption is one layer of defense.
Customers still need network security, access controls, monitoring.""")
    
    # =====================================================================
    # SLIDE 20: Competitive Landscape
    # =====================================================================
    add_content_slide(prs, "How We Differentiate", [
        "vs Oracle TDE:",
        "  - TDE decrypts in memory - DBAs see plaintext",
        "  - MongoDB: Server NEVER sees plaintext",
        "vs Cosmos DB:",
        "  - Cosmos is deterministic-only (pattern leakage)",
        "  - No range queries on encrypted data",
        "  - MongoDB QE: Random encryption + range queries",
        "vs PostgreSQL:",
        "  - Server-side encryption - decrypts for processing",
        "  - No built-in KMS integration"
    ], """Key message: MongoDB is the only document database with client-side 
searchable encryption that supports range queries while maintaining 
true zero-trust (server never sees plaintext).""")
    
    # =====================================================================
    # SLIDE 21: Discovery Questions
    # =====================================================================
    add_two_column_slide(prs, "Customer Discovery Questions",
        "Security Questions",
        [
            "What data classifications do you have?",
            "Who should NOT see sensitive data?",
            "Do you have insider threat concerns?",
            "What's your key management strategy?"
        ],
        "Compliance Questions",
        [
            "Which regulations apply?",
            "Have you had any audit findings?",
            "Do you need Right to be Forgotten?",
            "What are your data residency requirements?"
        ],
        """The killer question:
\"Should your database administrators be able to see customer SSNs?\"

The answer is almost always \"No\" - and that's where client-side 
encryption becomes essential. TDE doesn't help because DBAs can 
still see decrypted data.""")
    
    # =====================================================================
    # SLIDE 22: Objection Handling
    # =====================================================================
    add_content_slide(prs, "Common Objections & Responses", [
        "\"We already have TDE\" → TDE decrypts in memory. DBAs see plaintext.",
        "\"This will slow down queries\" → ~10% overhead. Security vs speed trade-off.",
        "\"We can't search encrypted data\" → QE enables equality AND range queries.",
        "\"Too complex to implement\" → Automatic encryption uses schema definition.",
        "\"Compliance is OK with TDE\" → Ask about insider threats. GDPR covers internal access."
    ], """The key: Don't argue features. 
Understand their specific concerns and address those.

Frame it as: \"What's the cost of a data breach vs. 10% latency?\" """)
    
    # =====================================================================
    # SLIDE 23: When NOT to Use
    # =====================================================================
    add_two_column_slide(prs, "Anti-Patterns & Bad Fits",
        "Don't Use When...",
        [
            "Heavy aggregations on encrypted fields",
            "Full-text search on sensitive data",
            "Sorting is critical on encrypted fields",
            "Sub-millisecond latency requirements",
            "Legacy apps that can't be modified"
        ],
        "Consider Instead...",
        [
            "Client-side aggregation for analytics",
            "Tokenized search with separate index",
            "Application-layer sorting after decrypt",
            "Accept latency for security trade-off",
            "Incremental migration path"
        ],
        """Be honest about when NOT to use CSFLE/QE. This builds trust.

The key: Qualify early. A failed POC is worse than a declined engagement.""")
    
    # =====================================================================
    # SLIDE 24: Labs Overview
    # =====================================================================
    add_content_slide(prs, "Three 34-Minute Labs", [
        "Lab 1: CSFLE & Troubleshooting (34 min)",
        "  - AWS KMS setup, Automatic encryption",
        "  - crypt_shared debugging, BSON Subtype 6 verification",
        "",
        "Lab 2: QE Range Queries (34 min)",
        "  - Range query on salary field",
        "  - Inspect .esc/.ecoc, DEK per field verification",
        "",
        "Lab 3: Right to Erasure (34 min)",
        "  - 1 DEK per user pattern",
        "  - Crypto-shredding demo, selective erasure verification"
    ], """Each lab has step-by-step instructions and checkpoints.
Labs use AWS KMS for realistic enterprise scenarios.""")
    
    # =====================================================================
    # SLIDE 25: Key Takeaways & Resources
    # =====================================================================
    add_content_slide(prs, "Key Takeaways", [
        "CSFLE = deterministic + equality queries (MongoDB 4.2+)",
        "QE = randomized + range queries (MongoDB 7.0+)",
        "QE requires separate DEK per field (critical difference)",
        "2-3x storage for Range indexes - plan capacity accordingly",
        ".esc and .ecoc need periodic compaction (monthly)",
        "1 DEK per user = crypto-shredding for GDPR compliance",
        "",
        "Resources: docs.mongodb.com/csfle, docs.mongodb.com/qe"
    ], """Final reminders:
1. CSFLE vs QE choice depends on query requirements
2. QE is more secure (randomized) but has overhead
3. The DEK-per-field requirement for QE is critical
4. Storage planning must account for 2-3x overhead
5. Compaction is an operational requirement
6. Crypto-shredding enables GDPR compliance""")
    
    return prs

def main():
    """Main function to generate the presentation."""
    print("Generating MongoDB CSFLE & QE SA Enablement Presentation...")
    
    prs = create_presentation()
    
    # Save the presentation
    output_path = os.path.join(os.path.dirname(__file__), 
                               "MongoDB_CSFLE_QE_SA_Enablement.pptx")
    prs.save(output_path)
    
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    main()
