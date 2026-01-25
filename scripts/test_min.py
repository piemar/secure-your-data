#!/usr/bin/env python3
from pptx import Presentation
import os

def main():
    prs = Presentation()
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Slide {i+1}"
        slide.placeholders[1].text = f"Content for slide {i+1}"
    
    output = "test_min.pptx"
    prs.save(output)
    print(f"DONE: {output}")

if __name__ == "__main__":
    main()
